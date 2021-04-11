# basic
import numpy as np
import pandas as pd

# get data
import pandas_datareader as pdr

# visual
import matplotlib.pyplot as plt
import mpl_finance as mpf
import seaborn as sns

#time
import datetime as datetime

#talib
import talib

#csv
import csv

#ptt
from pptx import Presentation 
from pptx.util import Inches

# ======= Create PPT file ==========
prs = Presentation('template.pptx')

# ======= select stock =========
# ------- volume_of_month_S&P500.csv file -------
# df = pd.read_csv('volume_of_month_S&P500.csv')		
# df["Company"] = df["Company"].str.replace(".","-")
# df.columns=["","Company","volume"]
# X = df[["Company"]]
# X = np.array(X)

# ------- stock -------
X = np.array(["AAL", "AAPL", "ACB", "AMD", "BA", "BAC", \
			  "BSX", "C", "CCL", "CMCSA", "CSCO", "DAL", \
			  "F", "FCX", "GE", "GM", "INTC", \
			  "MRO", "MSFT", "MU",  "NCLH", "OXY", \
			  "PFE", "T", "UAL", "WFC", "XOM"])
			  
# X = np.array(["DAL"])

# ======= Catch stock data ======= 
start = datetime.datetime(2020,4,27)
# end = datetime.datetime(2020,1,28)
end = datetime.datetime.now()

for i in range(27):	#len(X)
	#stock input format： 
	# TW：stock_number.TW(ex. 2330.TW) 
	# USA：stock_codename(ex.AAPL)	
	# company = str('GE')
	# df = pdr.DataReader(company, 'yahoo', start, end)
	company = "".join(X[i])	
	df = pdr.DataReader(company, 'yahoo', start, end)
	
	# ======= Use pandas to set index type ======= 
	# string formate time(strftime('%Y-%m-%d')) 
	df.index = df.index.format(formatter=lambda x: x.strftime('%m-%d')) 

	# ======= Create figure ======= 
	fig = plt.figure(figsize=(13, 7))
	ax = fig.add_subplot(1, 1, 1)

	# ======= Set x/y label interval ======= 
	ax.set_xticks(range(0, len(df.index), 10))
	ax.set_xticklabels(df.index[::10])

	# ======= Create stock k bar ======= 
	mpf.candlestick2_ochl(ax, df['Open'], df['Close'], df['High'],
						  df['Low'], width=0.6, colorup='r', colordown='g', alpha=0.75);

	# ======= Calculate average line	=======	
	sma_3 = talib.SMA(np.array(df['Close']), 3)	
	sma_5 = talib.SMA(np.array(df['Close']), 5)		  
	sma_20 = talib.SMA(np.array(df['Close']), 20)
	sma_60 = talib.SMA(np.array(df['Close']), 60)

	# ====== Calculate Bollinger bands ========
	upperband, middleband, lowerband = talib.BBANDS(df['Close'], timeperiod=20, \
													nbdevup=2, nbdevdn=2, matype=0) # matype = 0:SMA 1:EMA 2:WMA


	# ======= Set font type ==========
	plt.rcParams['font.sans-serif']=['Times New Roman'] 
	# plt.rcParams['font.sans-serif']=['Microsoft JhengHei'] # show chinese word

	# ======= Plot figure ==========
	# ax.plot(sma_3, label='3MA')
	ax.plot(sma_5, label='5MA')
	ax.plot(sma_20, label='20MA')
	# ax.plot(sma_60, label='60MA')
	ax.plot(upperband, label='BBand-up')
	ax.plot(lowerband, label='BBand-low')
	
	# ======= Set x/y axis limit ======= 
	# a = np.array(df['High'])
	# b = np.array(lowerband)
	# plt.xlim(len(df.index)-60, len(df.index))
	# plt.ylim(np.nanmin(b), np.max(a)+0.5)

	# print(df['Open'].tail())
	
	# ======== zoom_in_out_in_figure =========
	# import sys
	# sys.path.append('C:\\Users\\lankon\\Google 雲端硬碟 (s03212036@go.thu.edu.tw)\\股票分析\\python_code\\python_subroutine')
	# from matplotlib.pyplot import figure, show
	# from zoom_in_out_in_figure import ZoomPan
	# scale = 1.1 
	# zp = ZoomPan() 
	# figZoom = zp.zoom_factory(ax, base_scale = scale) 
	# figPan = zp.pan_factory(ax) 
		
	# ======== Plot figure ========
	ax.legend();
	plt.title(company, fontsize = 25)
	fig.savefig('plot.png')
	# plt.show()
	
	# ======= Create PPT file ==========		
	blank_slide_layout = prs.slide_layouts[i+6] 
	slide = prs.slides.add_slide(blank_slide_layout)
	# add_picture(image_file, left, top, width, height)
	slide.shapes.add_picture('plot.png', Inches(0), Inches(1), Inches(10), Inches(5.3)) 

# ======= Save PPT file ==========
end = end.strftime('%Y-%m-%d')
prs.save('C:\\Users\\lankon\\Google 雲端硬碟 (s03212036@go.thu.edu.tw)\\股票分析\\' +str(end) + '.pptx')				  